#!/usr/bin/env python3
"""
BFC Game Upload Script
Usage:
  python3 scripts/upload_game.py game.zip  folder_name [--push]
  python3 scripts/upload_game.py game.pptx folder_name [--push]

Supports two PPTX structures:

FORMAT A (Classic МыслиБой ZIP):
  - First half: org slides + question slides (each ends with Q number)
  - Second half: "ОТВЕТЫ" + answer slides (correct answer in FFC000 color)
  - Media files named by question number: 2.mp3, 5.mp4, etc.

FORMAT B (BFC Tournament triplets):
  - Every question is 3 slides: answers-text, question-image, answer-image
  - Org slides between rounds
  - Media files: q02_audio.mp3, q08_video.mp4, etc.
"""

import sys, os, re, json, shutil, subprocess, time, zipfile, tempfile
from pathlib import Path
from pptx import Presentation
from lxml import etree

REPO_ROOT = Path(__file__).parent.parent
MEDIA_ROOT = REPO_ROOT / "media"
SCRATCHPAD = Path("/tmp/bfc_upload")

GITHUB_BASE = "https://raw.githubusercontent.com/pantandrej/mind-fight-club/main/media/"
CORRECT_COLOR = "FFC000"  # gold highlight = correct answer


# ── Slide text extraction ───────────────────────────────────────

def slide_texts(slide) -> list[str]:
    """All non-empty text lines from a slide."""
    out = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    out.append(t)
    return out


def slide_colored_texts(slide) -> list[tuple[str, list[str]]]:
    """[(text, [colors])] for each paragraph in slide."""
    out = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            txt = para.text.strip()
            if not txt:
                continue
            xml = etree.tostring(para._p, pretty_print=False).decode()
            clrs = re.findall(r'srgbClr val="([^"]+)"', xml, re.IGNORECASE)
            out.append((txt, [c.upper() for c in clrs]))
    return out


# ── Format A detection ──────────────────────────────────────────

def _trailing_qnum(texts: list[str]) -> int | None:
    """If last text is a standalone integer 1-20, return it, else None."""
    if not texts:
        return None
    last = texts[-1].strip()
    try:
        n = int(last)
        if 1 <= n <= 30:
            return n
    except ValueError:
        pass
    return None


def detect_format_a(prs: Presentation) -> bool:
    """
    Format A: has an 'ОТВЕТЫ' slide roughly in the middle.
    """
    for slide in prs.slides:
        texts = [t.upper() for t in slide_texts(slide)]
        if any('ОТВЕТЫ' in t or 'ОТВЕТ' == t for t in texts):
            return True
    return False


def parse_format_a(prs: Presentation) -> dict:
    """Parse МыслиБой-style PPTX (questions first half, answers second half)."""
    total = len(prs.slides)
    print(f"📊 Format A detected. Total slides: {total}")

    # Find the ОТВЕТЫ divider slide index
    divider_idx = None
    for i, slide in enumerate(prs.slides):
        texts = [t.upper() for t in slide_texts(slide)]
        if any('ОТВЕТЫ' in t for t in texts):
            divider_idx = i
            break

    if divider_idx is None:
        raise ValueError("Could not find ОТВЕТЫ slide")

    print(f"   Divider at slide {divider_idx + 1}")

    # Question slides: slides before divider that have a trailing question number
    q_slides = {}  # qnum → slide_index
    org_slides = []
    for i in range(divider_idx):
        slide = prs.slides[i]
        texts = slide_texts(slide)
        qnum = _trailing_qnum(texts)
        if qnum is not None:
            q_slides[qnum] = i
        else:
            org_slides.append({"slide": i + 1, "texts": texts})

    # Answer slides: slides after divider (skip divider itself), match by qnum
    a_slides = {}  # qnum → slide_index
    for i in range(divider_idx + 1, total):
        slide = prs.slides[i]
        texts = slide_texts(slide)
        qnum = _trailing_qnum(texts)
        if qnum is not None:
            a_slides[qnum] = i

    print(f"   Q slides: {sorted(q_slides.keys())}")
    print(f"   A slides: {sorted(a_slides.keys())}")

    # Build questions
    questions = []
    for qnum in sorted(q_slides.keys()):
        q_idx = q_slides[qnum]
        a_idx = a_slides.get(qnum)

        q_texts = slide_texts(prs.slides[q_idx])
        # Answers: all lines except trailing number and non-option lines
        answers = _extract_options(q_texts)
        question_text = _extract_question_text(q_texts, answers)

        # Correct answer: FFC000-highlighted option on answer slide
        correct = None
        if a_idx is not None:
            colored = slide_colored_texts(prs.slides[a_idx])
            a_texts = slide_texts(prs.slides[a_idx])
            a_options = _extract_options(a_texts)
            for txt, colors in colored:
                if CORRECT_COLOR in colors and txt in a_options:
                    # Map back to question option index
                    # Options may be the same text in same order
                    for ci, opt in enumerate(answers):
                        if opt == txt:
                            correct = ci
                            break
                    if correct is None:
                        # Try partial match
                        for ci, opt in enumerate(answers):
                            if txt in opt or opt in txt:
                                correct = ci
                                break

        questions.append({
            "n": qnum,
            "question_text": question_text,
            "answers": answers,
            "correct": correct,
            "slides": {"question": q_idx + 1, "answer": a_idx + 1 if a_idx is not None else None},
            "slide_q_url": None,
            "slide_a_url": None,
            "audio": None, "video": None,
            "answer_audio": None, "answer_video": None,
        })

    # End slide (after all answer slides)
    for i in range(divider_idx + 1, total):
        texts = slide_texts(prs.slides[i])
        if _trailing_qnum(texts) is None and i > max(a_slides.values(), default=0):
            org_slides.append({"slide": i + 1, "texts": texts})

    return {"format": "A", "org_slides": org_slides, "questions": questions, "total": total,
            "divider": divider_idx + 1}


def _extract_options(texts: list[str]) -> list[str]:
    """
    Extract answer options from slide texts.
    Options are lines that are NOT the question text (not ending with '?')
    and NOT standalone numbers (question numbers).
    Heuristic: options come AFTER the question line.
    """
    if not texts:
        return []

    # Find question line (usually long or ends with ?)
    q_idx = 0
    for i, t in enumerate(texts):
        if len(t) > 20 or t.endswith('?'):
            q_idx = i
            break

    # The LAST standalone integer in the slide is the question number — skip only that.
    # All other numbers can be valid answer options (e.g. "1","2","3","4","5").
    last_int_idx = None
    for i, t in enumerate(texts):
        try:
            int(t.strip())
            last_int_idx = i
        except ValueError:
            pass

    options = []
    for i, t in enumerate(texts[q_idx + 1:], start=q_idx + 1):
        if i == last_int_idx:
            continue  # skip question number
        options.append(t)

    # If nothing found, include all lines except question text and last number
    if not options:
        for i, t in enumerate(texts):
            if i == last_int_idx:
                continue
            if not (len(t) > 25 or t.endswith('?')):
                options.append(t)

    # If mix of multi-char and single-digit — single digits are layout artifacts, drop them.
    # Exception: if ALL options are single digits (e.g. Q10 "which Fabrika Zvezd: 1,2,3,4,5")
    multi = [o for o in options if not (o.strip().lstrip('-').isdigit() and len(o.strip()) <= 2)]
    single_only = [o for o in options if o.strip().lstrip('-').isdigit() and len(o.strip()) <= 2]
    if multi and single_only:
        options = multi  # drop single digits — they're visual slide labels, not answer options

    return options[:8]  # cap at 8 to avoid runaway cases


def _extract_question_text(texts: list[str], answers: list[str]) -> str:
    """Extract question text (not in answers, not the trailing question number)."""
    ans_set = set(answers)
    # Last standalone int = question number, skip it
    last_int_idx = None
    for i, t in enumerate(texts):
        try:
            int(t.strip())
            last_int_idx = i
        except ValueError:
            pass
    q_lines = []
    for i, t in enumerate(texts):
        if i == last_int_idx:
            continue
        if t not in ans_set:
            q_lines.append(t)
    return ' '.join(q_lines).strip()


# ── Format B detection ──────────────────────────────────────────

def parse_format_b(prs: Presentation) -> dict:
    """Parse BFC triplet-style PPTX."""
    total = len(prs.slides)
    print(f"📊 Format B detected. Total slides: {total}")

    # Answer slides have numbered options (2+ numbered lines)
    def is_answer_slide(slide):
        texts = slide_texts(slide)
        joined = '\n'.join(texts)
        num = len(re.findall(r'(?:^|\n)\d+[.)]\s', joined))
        let = len(re.findall(r'(?:^|\n)[А-ЯA-Z][.)]\s', joined))
        return num >= 2 or let >= 2

    org_slides = []
    questions = []
    i = 0
    while i < total:
        slide = prs.slides[i]
        if is_answer_slide(slide):
            qnum = len(questions) + 1
            texts = slide_texts(slide)
            from scripts_utils_placeholder import parse_answers_from_texts
            q_text, answers = parse_answers_from_texts(texts)
            questions.append({
                "n": qnum,
                "question_text": q_text,
                "answers": answers,
                "correct": None,
                "slides": {"answers": i + 1, "question": i + 2, "answer": i + 3},
                "slide_q_url": None, "slide_a_url": None,
                "audio": None, "video": None,
                "answer_audio": None, "answer_video": None,
            })
            i += 3
        else:
            org_slides.append({"slide": i + 1, "texts": slide_texts(slide)})
            i += 1

    return {"format": "B", "org_slides": org_slides, "questions": questions, "total": total}


# ── Inline Format B option parser (no import needed) ───────────

def parse_answers_from_texts(texts: list[str]) -> tuple[str, list[str]]:
    answer_pattern = re.compile(r'^[\d]+[.)]\s+(.+)$|^[А-ЯA-Z][.)]\s+(.+)$')
    answers, q_lines = [], []
    for text in texts:
        m = answer_pattern.match(text)
        if m:
            answers.append((m.group(1) or m.group(2)).strip())
        else:
            q_lines.append(text)
    return ' '.join(q_lines).strip(), answers


# ── Keynote slide export ────────────────────────────────────────

def export_slides_pdf(pdf_path: Path, out_dir: Path):
    """Render PDF pages to JPEG using PyMuPDF (no Keynote needed)."""
    import fitz  # PyMuPDF
    out_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(pdf_path))
    print(f"⏳ Rendering {len(doc)} slides from PDF...")
    for i, page in enumerate(doc, 1):
        mat = fitz.Matrix(2.0, 2.0)  # 2x = ~144dpi, good quality
        pix = page.get_pixmap(matrix=mat)
        out_path = out_dir / f"Slide.{i:03d}.jpg"
        pix.save(str(out_path))
    doc.close()
    print(f"✅ Rendered {len(list(out_dir.glob('*.jpg')))} slides")


def build_slide_map(slide_dir: Path) -> dict[int, Path]:
    """Slide number → JPEG path. Handles both Keynote and PyMuPDF naming."""
    result = {}
    for f in slide_dir.iterdir():
        if f.suffix.lower() not in ('.jpg', '.jpeg'):
            continue
        # Extract any number from filename
        digits = re.findall(r'\d+', f.stem)
        if digits:
            result[int(digits[-1])] = f
    return result


# ── Media matching ──────────────────────────────────────────────

def match_media_format_a(pptx_dir: Path, base_url: str, questions: list[dict]):
    """Format A: media files named by question number (2.mp3, 5.mp4)."""
    for f in pptx_dir.iterdir():
        if f.suffix.lower() not in ('.mp3', '.mp4', '.wav', '.m4a'):
            continue
        try:
            qnum = int(f.stem)
        except ValueError:
            continue
        for q in questions:
            if q["n"] == qnum:
                url = base_url + f.name
                if f.suffix.lower() in ('.mp3', '.wav', '.m4a'):
                    q["audio"] = url
                else:
                    q["video"] = url
                print(f"  🎵 Q{qnum} media: {f.name}")


def match_media_format_b(pptx_dir: Path, base_url: str, questions: list[dict]):
    """Format B: media files named q02_audio.mp3, q08_video.mp4, etc."""
    for f in pptx_dir.iterdir():
        if f.suffix.lower() not in ('.mp3', '.mp4', '.wav', '.m4a'):
            continue
        m = re.match(r'q(\d+)_?(audio|video|a_video|a_audio)', f.stem, re.I)
        if not m:
            continue
        qn, kind = int(m.group(1)), m.group(2).lower()
        url = base_url + f.name
        for q in questions:
            if q["n"] == qn:
                if kind == 'audio':       q["audio"] = url
                elif kind == 'video':     q["video"] = url
                elif kind == 'a_video':   q["answer_video"] = url
                elif kind == 'a_audio':   q["answer_audio"] = url


# ── Main ────────────────────────────────────────────────────────

def main():
    if len(sys.argv) < 3:
        print("Usage: python3 scripts/upload_game.py <file.pptx|file.zip> <folder> [--push]")
        print("Example: python3 scripts/upload_game.py 'Мысли Бой 3.zip' mb3 --push")
        sys.exit(1)

    input_path = Path(sys.argv[1]).expanduser().resolve()
    folder = sys.argv[2].strip('/')
    do_push = '--push' in sys.argv

    if not input_path.exists():
        print(f"❌ Not found: {input_path}")
        sys.exit(1)

    # Extract ZIP if needed
    work_dir = SCRATCHPAD / folder / "source"
    work_dir.mkdir(parents=True, exist_ok=True)

    if input_path.suffix.lower() == '.zip':
        print(f"📦 Extracting ZIP...")
        with zipfile.ZipFile(input_path) as zf:
            zf.extractall(work_dir)
        # Find PPTX inside
        pptx_files = list(work_dir.rglob("*.pptx"))
        pptx_files = [f for f in pptx_files if not f.name.startswith('._')]
        if not pptx_files:
            print("❌ No .pptx found in ZIP")
            sys.exit(1)
        pptx_path = pptx_files[0]
        pptx_dir = pptx_path.parent
        print(f"   Found: {pptx_path.name}")
    else:
        pptx_path = input_path
        pptx_dir = pptx_path.parent

    # Export slides — prefer PDF (fast, no Keynote needed)
    slide_out = SCRATCHPAD / folder / "slides"
    pdf_files = [f for f in pptx_dir.rglob("*.pdf") if not f.name.startswith('._')]
    if pdf_files:
        print(f"📄 Found PDF: {pdf_files[0].name} — using it instead of Keynote")
        export_slides_pdf(pdf_files[0], slide_out)
    else:
        print("⚠️  No PDF found — falling back to Keynote export")
        export_slides_keynote_fallback(pptx_path, slide_out)
    slide_map = build_slide_map(slide_out)
    print(f"   Slide map: {len(slide_map)} slides ({min(slide_map)} - {max(slide_map)})")

    # Parse PPTX
    prs = Presentation(str(pptx_path))
    if detect_format_a(prs):
        structure = parse_format_a(prs)
    else:
        structure = parse_format_b(prs)

    fmt = structure["format"]
    dest_dir = MEDIA_ROOT / folder
    dest_dir.mkdir(parents=True, exist_ok=True)
    base_url = GITHUB_BASE + folder + "/"

    # Copy org slide images
    org_out = []
    for idx, org in enumerate(structure["org_slides"]):
        snum = org["slide"]
        name = f"org_{idx+1:02d}.jpg"
        src = slide_map.get(snum)
        if src:
            shutil.copy(src, dest_dir / name)
            org["url"] = base_url + name
            org["filename"] = name
            print(f"  📄 Org slide {snum} → {name}")
        org_out.append(org)

    # Copy question / answer slide images
    for q in structure["questions"]:
        n = q["n"]
        sq_src = slide_map.get(q["slides"].get("question"))
        sa_src = slide_map.get(q["slides"].get("answer"))
        sq_name = f"q{n:02d}_sq.jpg"
        sa_name = f"q{n:02d}_sa.jpg"
        if sq_src:
            shutil.copy(sq_src, dest_dir / sq_name)
            q["slide_q_url"] = base_url + sq_name
            print(f"  🖼  Q{n} question → {sq_name}")
        if sa_src:
            shutil.copy(sa_src, dest_dir / sa_name)
            q["slide_a_url"] = base_url + sa_name
            print(f"  🖼  Q{n} answer  → {sa_name}")

    # Copy media and match to questions
    for f in pptx_dir.iterdir():
        if f.suffix.lower() in ('.mp3', '.mp4', '.wav', '.m4a') and not f.name.startswith('._'):
            shutil.copy(f, dest_dir / f.name)
    if fmt == "A":
        match_media_format_a(pptx_dir, base_url, structure["questions"])
    else:
        match_media_format_b(pptx_dir, base_url, structure["questions"])

    # Auto-detected summary
    auto_correct = sum(1 for q in structure["questions"] if q["correct"] is not None)
    print(f"\n✅ Auto-detected correct answers: {auto_correct}/{len(structure['questions'])}")
    for q in structure["questions"]:
        if q["correct"] is not None:
            print(f"   Q{q['n']}: [{q['correct']+1}] {q['answers'][q['correct']]}")
        else:
            print(f"   Q{q['n']}: ❓ нужно выбрать вручную")

    # Write game.json
    game_json = {
        "folder": folder,
        "format": fmt,
        "base_url": base_url,
        "total_slides": structure["total"],
        "org_slides": org_out,
        "questions": structure["questions"],
        "published": False,
    }
    json_path = dest_dir / "game.json"
    with open(json_path, 'w', encoding='utf-8') as f:
        json.dump(game_json, f, ensure_ascii=False, indent=2)
    print(f"\n📄 game.json → {json_path}")

    # Git commit + push
    print("\n📤 Committing to GitHub...")
    subprocess.run(["git", "add", f"media/{folder}/"], cwd=REPO_ROOT, check=True)
    subprocess.run(["git", "commit", "-m",
                    f"media: upload game '{folder}' ({len(structure['questions'])} questions, format {fmt})"],
                   cwd=REPO_ROOT, check=True)
    if do_push:
        subprocess.run(["git", "push", "origin", "main"], cwd=REPO_ROOT, check=True)
        print("✅ Pushed!")
    else:
        print("ℹ️  git push origin main — to push")

    print(f"\n🎮 Открой Конструктор и загрузи:")
    print(f"   {base_url}game.json")


if __name__ == "__main__":
    main()
