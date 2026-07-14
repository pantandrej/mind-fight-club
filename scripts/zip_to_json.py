#!/usr/bin/env python3
"""
BFC Game Importer
Usage: python3 scripts/zip_to_json.py /path/to/game.zip --code T1 [--github-repo pantandrej/mind-fight-club]

Exports PPTX slides as images, uploads everything to Supabase Storage,
parses questions/answers from slide text, and generates a ready-to-import
JSON for the admin panel (Импорт игр).
"""
import sys, os, json, re, subprocess, tempfile, zipfile, requests, argparse, shutil
from pathlib import Path

env_path = Path(__file__).parent / '.env'
creds = {}
if env_path.exists():
    for line in env_path.read_text().splitlines():
        if '=' in line and not line.startswith('#'):
            k, v = line.split('=', 1)
            creds[k.strip()] = v.strip()

SUPABASE_URL = creds.get('SUPABASE_URL', '')
SERVICE_KEY  = creds.get('SUPABASE_SERVICE_KEY', '')
GITHUB_REPO  = 'pantandrej/mind-fight-club'

HEADERS = {
    "apikey": SERVICE_KEY,
    "Authorization": f"Bearer {SERVICE_KEY}",
}

def upload_to_storage(local_path, storage_path):
    mime = 'audio/mpeg' if local_path.endswith('.mp3') else \
           'video/mp4'  if local_path.endswith('.mp4') else \
           'image/jpeg' if local_path.endswith(('.jpg','.jpeg')) else 'application/octet-stream'
    with open(local_path, 'rb') as f:
        r = requests.post(
            f"{SUPABASE_URL}/storage/v1/object/mfc-media/{storage_path}",
            headers={**HEADERS, "Content-Type": mime, "x-upsert": "true"},
            data=f
        )
    if not r.ok:
        print(f"  ⚠️  Upload failed {storage_path}: {r.status_code} {r.text[:100]}")
        return None
    return f"{SUPABASE_URL}/storage/v1/object/public/mfc-media/{storage_path}"

def export_slides_keynote(pptx_path, out_dir):
    script = f'''
tell application "Keynote"
    set pptxPath to POSIX file "{pptx_path}"
    set outFolder to POSIX file "{out_dir}/"
    open pptxPath
    delay 3
    set theDoc to front document
    export theDoc to outFolder as slide images with properties {{image format:JPEG, compression factor:0.85}}
    close theDoc saving no
end tell
'''
    result = subprocess.run(['osascript', '-e', script], capture_output=True, text=True)
    return result.returncode == 0

def parse_pptx_slides(pptx_path):
    """Returns list of slide text lists."""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        result = []
        for slide in prs.slides:
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            result.append(texts)
        return result
    except Exception as e:
        print(f"  ⚠️  PPTX parse error: {e}")
        return []

# Patterns for organizational/info slides
INFO_PATTERNS = re.compile(
    r'^(ПЕРВЫЙ|ВТОРОЙ|ТРЕТИЙ|ЧЕТВЁРТЫЙ|ПЯТЫЙ|ШЕСТОЙ|СЕДЬМОЙ|ВОСЬМОЙ|'
    r'первый|второй|третий|РАУНД|раунд|СПАСИБО|ПРАВИЛА|правила|'
    r'ПИШИТЕ|СКОЛЬКО БАЛЛОВ|ИТОГО|ФИНАЛ|РАЗМИНКА|РАЗОГРЕВ)',
    re.IGNORECASE
)

def classify_slide(texts, slide_num):
    """Returns 'info', 'question', 'answer_only', or 'skip'."""
    if not texts:
        return 'skip'
    full = ' '.join(texts)
    first = texts[0] if texts else ''

    # Short single-text slides (round headers etc)
    if len(texts) <= 2 and len(full) < 60:
        if INFO_PATTERNS.match(first):
            return 'info'
        # Slides with just a number are answer-only
        if re.match(r'^\d+$', full.strip()):
            return 'answer_only'

    # Slides that look like "answer only" (question text + number, no answer options)
    has_numbered_options = bool(re.search(r'\b[1-6]\.\s+\S', full))
    if not has_numbered_options and len(texts) <= 3:
        return 'answer_only'

    if has_numbered_options:
        return 'question'

    # Info patterns anywhere in text
    if INFO_PATTERNS.match(first):
        return 'info'

    return 'question'

def extract_answers_from_texts(texts):
    """Extract numbered answer options like '1. Foo', '2. Bar'."""
    answers = []
    for t in texts:
        m = re.match(r'^[1-6]\.\s+(.+)', t.strip())
        if m:
            answers.append(m.group(1).strip())
    return answers

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('zip_path')
    parser.add_argument('--code', required=True, help='Game code slug (e.g. t1)')
    parser.add_argument('--github-repo', default=GITHUB_REPO)
    args = parser.parse_args()

    code = args.code.lower()

    print(f"🎮 Processing → code={code}")

    with tempfile.TemporaryDirectory() as tmpdir:
        # Extract zip
        with zipfile.ZipFile(args.zip_path) as z:
            z.extractall(tmpdir)

        pptx_files = list(Path(tmpdir).rglob('*.pptx'))
        if not pptx_files:
            print("❌ No PPTX found"); sys.exit(1)
        pptx_path = str(pptx_files[0])
        print(f"  📊 PPTX: {Path(pptx_path).name}")

        # Find media (skip files with о/o suffix — those are answer videos)
        media_files = {}
        answer_media = {}
        for ext in ['*.mp3', '*.mp4', '*.wav']:
            for f in Path(tmpdir).rglob(ext):
                if '__MACOSX' in str(f): continue
                name = f.name
                # Files ending in о (cyrillic) or _answer are answer reveal media
                stem = f.stem
                if stem.endswith('о') or stem.endswith('o') and stem[:-1].isdigit():
                    answer_media[name] = str(f)
                else:
                    media_files[name] = str(f)

        print(f"  🎵 Question media: {list(media_files.keys())}")
        print(f"  🎵 Answer media: {list(answer_media.keys())}")

        # Export slides via Keynote
        slides_dir = os.path.join(tmpdir, 'slides')
        os.makedirs(slides_dir, exist_ok=True)
        print("🖼  Exporting slides via Keynote...")
        if not export_slides_keynote(pptx_path, slides_dir):
            print("❌ Keynote export failed"); sys.exit(1)

        slide_files = sorted([f for f in os.listdir(slides_dir) if f.endswith(('.jpeg','.jpg'))])
        n_slides = len(slide_files)
        print(f"  ✅ {n_slides} slides")

        # Upload slides to Supabase Storage
        print("☁️  Uploading slides...")
        slide_urls = {}
        for i, fname in enumerate(slide_files, 1):
            local = os.path.join(slides_dir, fname)
            url = upload_to_storage(local, f"games/{code}/slide_{i:03d}.jpg")
            if url:
                slide_urls[i] = url
                print(f"  ✅ slide {i:03d}")
            else:
                # Fallback: GitHub raw URL (if slides already pushed)
                gh = f"https://raw.githubusercontent.com/{args.github_repo}/main/public/games/{code}/slide_{i:03d}.jpg"
                slide_urls[i] = gh
                print(f"  ↩️  slide {i:03d} → GitHub fallback")

        # Upload media
        print("☁️  Uploading media...")
        media_urls = {}
        for fname, local in {**media_files, **answer_media}.items():
            safe_name = re.sub(r'[^\w.\-]', '_', fname)
            url = upload_to_storage(local, f"games/{code}/media/{safe_name}")
            if url:
                media_urls[fname] = url
                print(f"  ✅ {fname}")

        # Parse slide texts
        slide_texts = parse_pptx_slides(pptx_path)
        while len(slide_texts) < n_slides:
            slide_texts.append([])

        # Build question list: group slides into (question_slide, answer_slide) pairs
        # Pattern: info | (question_slide, skip_slide?, answer_slide) repeating
        questions = []
        i = 0
        while i < n_slides:
            slide_num = i + 1
            texts = slide_texts[i] if i < len(slide_texts) else []
            kind = classify_slide(texts, slide_num)

            if kind == 'info':
                questions.append({
                    "question": texts[0] if texts else f"Слайд {slide_num}",
                    "answer_a": "", "answer_b": "", "answer_c": "",
                    "answer_d": "", "answer_e": "", "answer_f": "",
                    "correct_answer": "A",
                    "category": "GENERAL",
                    "explanation": "",
                    "question_type": "info",
                    "slide_img_url": slide_urls.get(slide_num),
                    "answer_slide_img_url": None,
                    "audio_url": None, "video_url": None,
                })
                i += 1

            elif kind == 'question':
                q_slide = slide_num
                q_texts = texts
                answers = extract_answers_from_texts(q_texts)
                q_text = q_texts[0] if q_texts else f"Вопрос {slide_num}"

                # Next slide: skip (question without answers) — skip it
                # Slide after that: answer slide
                a_slide_num = None
                # Look ahead for answer slide (skip intermediate)
                if i + 2 < n_slides:
                    next_kind = classify_slide(slide_texts[i+1] if i+1 < len(slide_texts) else [], i+2)
                    after_kind = classify_slide(slide_texts[i+2] if i+2 < len(slide_texts) else [], i+3)
                    if next_kind in ('answer_only', 'skip') and after_kind in ('answer_only', 'skip', 'question', 'info'):
                        a_slide_num = i + 3  # 3rd slide is answer reveal
                        i += 3
                    else:
                        i += 1
                else:
                    i += 1

                # Determine media for this question (by question number in text)
                q_num_match = re.search(r'\b(\d+)\b', ' '.join(q_texts[-2:]))
                q_num = int(q_num_match.group(1)) if q_num_match else None

                audio_url = None
                video_url = None
                for fname, url in media_urls.items():
                    stem = Path(fname).stem
                    if q_num and (stem == str(q_num) or stem.rstrip('о') == str(q_num)):
                        if fname.endswith('.mp3'):
                            audio_url = url
                        elif fname.endswith('.mp4') and not stem.endswith('о'):
                            video_url = url

                # Map answers to fields
                abc = ['a','b','c','d','e','f']
                entry = {
                    "question": q_text,
                    "answer_a": answers[0] if len(answers) > 0 else "",
                    "answer_b": answers[1] if len(answers) > 1 else "",
                    "answer_c": answers[2] if len(answers) > 2 else "",
                    "answer_d": answers[3] if len(answers) > 3 else "",
                    "answer_e": answers[4] if len(answers) > 4 else "",
                    "answer_f": answers[5] if len(answers) > 5 else "",
                    "correct_answer": "A",  # user sets this in admin panel
                    "category": "GENERAL",
                    "explanation": "",
                    "question_type": "multiple_choice",
                    "slide_img_url": slide_urls.get(q_slide),
                    "answer_slide_img_url": slide_urls.get(a_slide_num) if a_slide_num else None,
                    "audio_url": audio_url,
                    "video_url": video_url,
                }
                questions.append(entry)

            else:  # answer_only or skip
                i += 1

        out_path = Path(__file__).parent / f"{code}_import.json"
        out_path.write_text(json.dumps(questions, ensure_ascii=False, indent=2), encoding='utf-8')

        valid = sum(1 for q in questions if q.get('question_type') != 'info')
        info  = sum(1 for q in questions if q.get('question_type') == 'info')
        print(f"\n✅ {out_path.name} готов")
        print(f"   Вопросов: {valid} | Орг.слайдов: {info}")
        print(f"\n👉 Загрузи в Админка → Импорт игр → выбери {out_path.name}")

if __name__ == '__main__':
    main()
